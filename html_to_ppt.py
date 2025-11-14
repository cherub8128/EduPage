"""
HTML 슬라이드를 PPT로 변환하는 스크립트
각 슬라이드를 캡처하여 PowerPoint 파일로 생성합니다.
"""

from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import time
import os
from pathlib import Path

def capture_slides_to_ppt(html_path, output_ppt_path, total_slides=49):
    """
    HTML 슬라이드를 캡처하여 PPT 파일로 생성

    Args:
        html_path: HTML 파일 경로
        output_ppt_path: 출력 PPT 파일 경로
        total_slides: 총 슬라이드 수
    """

    # 임시 이미지 저장 폴더
    temp_folder = Path("temp_slides")
    temp_folder.mkdir(exist_ok=True)

    print(f"🚀 슬라이드 캡처 시작: {html_path}")
    print(f"📊 총 {total_slides}개 슬라이드")

    # Playwright로 브라우저 실행 및 캡처
    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        page = browser.new_page(viewport={'width': 1920, 'height': 1080})

        # HTML 파일 열기
        file_url = f"file://{os.path.abspath(html_path).replace(os.sep,'/')}"
        page.goto(file_url)

        # KaTeX 수식 렌더링 대기
        time.sleep(2)

        # 각 슬라이드 캡처
        screenshot_paths = []

        for slide_num in range(1, total_slides + 1):
            print(f"📸 슬라이드 {slide_num}/{total_slides} 캡처 중...")

            # 슬라이드로 이동
            page.evaluate(f"showSlide({slide_num})")

            # 애니메이션 완료 대기
            time.sleep(0.5)

            # 스크린샷 저장
            screenshot_path = temp_folder / f"slide_{slide_num:03d}.png"
            page.screenshot(path=str(screenshot_path), full_page=False)
            screenshot_paths.append(screenshot_path)

        browser.close()

    print("\n📦 PPT 파일 생성 중...")

    # PowerPoint 프레젠테이션 생성
    prs = Presentation()
    prs.slide_width = Inches(16)  # 16:9 비율
    prs.slide_height = Inches(9)

    # 각 이미지를 슬라이드로 추가
    for i, img_path in enumerate(screenshot_paths, 1):
        print(f"➕ 슬라이드 {i}/{total_slides} 추가 중...")

        # 빈 슬라이드 추가
        blank_slide_layout = prs.slide_layouts[6]  # 빈 레이아웃
        slide = prs.slides.add_slide(blank_slide_layout)

        # 이미지 크기 조정 및 추가
        img = Image.open(img_path)
        img_width, img_height = img.size

        # 슬라이드 크기에 맞게 조정
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        # 비율 유지하면서 슬라이드에 맞추기
        img_ratio = img_width / img_height
        slide_ratio = slide_width / slide_height

        if img_ratio > slide_ratio:
            # 이미지가 더 넓음 - 너비에 맞춤
            pic_width = slide_width
            pic_height = int(slide_width / img_ratio)
        else:
            # 이미지가 더 높음 - 높이에 맞춤
            pic_height = slide_height
            pic_width = int(slide_height * img_ratio)

        # 중앙 정렬
        left = (slide_width - pic_width) // 2
        top = (slide_height - pic_height) // 2

        # 이미지 추가
        slide.shapes.add_picture(
            str(img_path),
            left, top,
            width=pic_width,
            height=pic_height
        )

    # PPT 파일 저장
    prs.save(output_ppt_path)
    print(f"\n✅ 완료! PPT 파일 생성: {output_ppt_path}")

    # 임시 파일 정리
    print("\n🧹 임시 파일 정리 중...")
    for img_path in screenshot_paths:
        img_path.unlink()
    temp_folder.rmdir()

    print("✨ 모든 작업 완료!")


if __name__ == "__main__":
    # 설정
    HTML_FILE = r"Teacher\MathEdu\TractrixPresentation.html"
    OUTPUT_PPT = "Tractrix_Presentation.pptx"
    TOTAL_SLIDES = 51

    # 실행
    capture_slides_to_ppt(HTML_FILE, OUTPUT_PPT, TOTAL_SLIDES)
